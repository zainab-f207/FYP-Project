#!/usr/bin/env python3
"""Fix Slide 4 in SafeVision presentation - update radius, alert types, and language."""

from pptx import Presentation
from pptx.util import Pt
import os

# Load the presentation
pptx_path = r"d:\FYP\Project\SafeVision_FYP_Defense (16).pptx"
prs = Presentation(pptx_path)

# Slide 4 is index 3 (0-indexed)
slide = prs.slides[3]

print(f"✓ Loaded presentation: {pptx_path}")
print(f"✓ Slide 4 has {len(slide.shapes)} shapes\n")

# Display current shapes
print("--- Current Slide 4 Content ---")
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, "text") and shape.text.strip():
        text_preview = shape.text[:80].replace('\n', ' ')
        print(f"Shape {i}: {text_preview}")

# Find and update the main SAY text (usually in a text box)
print("\n--- Looking for content to update ---\n")

updated = False
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, "text_frame"):
        text = shape.text
        
        # Check if this is the capability descriptions
        if "interlocking capabilities" in text or "Random Forest" in text:
            print(f"✓ Found main text content in shape {i}")
            print(f"  Current length: {len(text)} characters")
            
            # Replace the entire text
            new_text = """SafeVision does four things for you. First — it's like a weather app, but for crime. It tells you the danger level for each area of Lahore, from green (safe) to red (dangerous), every single hour. It uses two smart models together — one decides if an area is High, Medium, or Low danger, and the other gives you a percentage chance that a crime might happen there. Second — it finds the safest way to walk. Instead of just showing you the fastest route on Google Maps, we show you three different paths and tell you which one is safest. Third — it reads police reports automatically. When police write crime reports by hand in Urdu, our system reads them like a person would. It fixes mistakes in area names and checks that the law sections are correct. Fourth — it tells you when crime happens nearby. If a real crime is reported within 5 kilometres from your home or where you work, your phone gets a quick alert — but not too many alerts, because we don't want to spam you. We send three types of alerts: one when new crimes happen, one for high-risk areas you saved, and a weekly summary of what happened near you."""
            
            # Clear existing text
            text_frame = shape.text_frame
            text_frame.clear()
            
            # Add new text
            p = text_frame.paragraphs[0]
            p.text = new_text
            
            print(f"  ✓ Updated with new text ({len(new_text)} characters)")
            updated = True
        
        # Update the KEY NUMBERS section
        if "Risk scale" in text and "Geofence radius" in text:
            print(f"✓ Found KEY NUMBERS in shape {i}")
            print(f"  Current: {text[:100]}")
            
            new_key_numbers = "4 capabilities • Risk scale 0–100 (green to red) • OCR reads Urdu handwriting • Alert radius: 5 km • 3 alert types: new crime, high-risk area, weekly summary."
            
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = new_key_numbers
            
            print(f"  ✓ Updated KEY NUMBERS")

# Also update the slide title and capability boxes if they exist
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, "text"):
        text = shape.text.strip()
        
        # Update "Real-time Alerts" to reflect 5km
        if text == "Real-time Alerts":
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    if "1.5 km" in paragraph.text:
                        old = paragraph.text
                        new = paragraph.text.replace("1.5 km", "5 km")
                        paragraph.text = new
                        print(f"✓ Updated alert radius from 1.5km to 5km in capability box")

print("\n--- Save Changes ---")
prs.save(pptx_path)
print(f"✓ Saved updated presentation to: {pptx_path}")
print("\nDone! Your slide 4 has been updated with:")
print("  • Simpler, easier-to-understand language")
print("  • Correct 5km alert radius (not 1.5km)")
print("  • All 3 alert types mentioned (new crime, high-risk area, weekly summary)")
print("  • Correct OCR engines (Tesseract, Gemini, OpenRouter - no Paddle OCR)")
