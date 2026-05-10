#!/usr/bin/env python3
"""Fix remaining 1.5km references in slide 4."""

from pptx import Presentation

pptx_path = r"d:\FYP\Project\SafeVision_FYP_Defense (16).pptx"
prs = Presentation(pptx_path)
slide = prs.slides[3]  # Slide 4

print("✓ Checking and fixing all 1.5 km -> 5 km references...\n")

fixed_count = 0
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "1.5" in run.text and "km" in run.text:
                    old_text = run.text
                    run.text = run.text.replace("1.5 km", "5 km").replace("1.5km", "5km")
                    print(f"  Shape {i}: '{old_text}' → '{run.text}'")
                    fixed_count += 1

prs.save(pptx_path)
print(f"\n✓ Fixed {fixed_count} references")
print("✓ Presentation saved!")
